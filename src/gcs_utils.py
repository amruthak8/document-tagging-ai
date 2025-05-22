import os
import fitz  # PyMuPDF
from PIL import Image as Img
import traceback
from vertexai.generative_models import GenerativeModel,Image,Part, GenerationConfig
import vertexai
import json
import zipfile
import tempfile
import pandas as pd
from PIL import Image as Img
from pptx import Presentation
from google.cloud import storage
import config
import openpyxl 
from io import BytesIO
from docx import Document
import base64
import json
import xml.etree.ElementTree as ET
import fitz  
import csv
import io
from openpyxl import load_workbook
from datetime import datetime
from bs4 import BeautifulSoup
import time
from google.cloud import bigquery
from datetime import datetime

import logging


class GCSFileHandler:
    def __init__(self, bucket_name):
        """Initialize the GCS client and bucket."""
        self.client = storage.Client()
        self.bucket = self.client.bucket(bucket_name)
        self.bucket_name = bucket_name
    
    def create_folder(self, folder_name):
        
        """
        Creates a virtual folder in a Google Cloud Storage (GCS) bucket.
        
        This method creates a "folder" by uploading an empty object with a name ending in a slash (`/`),
        which most GCS interfaces interpret as a folder structure.

        Args:
            folder_name (str): The name of the folder to create. Can include nested paths like 'parent/child'.

        Returns:
            bool: True if the folder was successfully created.

        Logging:
            - Logs the start and end of folder creation.
            - Logs confirmation of the folder creation in the specified bucket.
        """
        
        logging.info("Printing in create_folder")
        
        """Creates a 'folder' in GCS by uploading an empty file with a trailing slash."""
        blob = self.bucket.blob(f"{folder_name}/")
        blob.upload_from_string("")  # Empty content
        
        logging.info(f"Folder '{folder_name}' created in bucket {self.bucket.name}.")
        
        logging.info("End of create_folder")
        
        return True

    
    def write_file_to_gcs(self, folder_name, file_name, content):
        
        """
        Writes (or overwrites) a file in a specified Google Cloud Storage (GCS) folder.

        This method uploads the provided content as a file into the designated folder within the GCS bucket.
        It supports writing all file types as long as the content is passed as a string (e.g., text, CSV, JSON, etc.).

        Args:
            folder_name (str): The target folder (path prefix) in the GCS bucket where the file will be stored.
            file_name (str): The name of the file to be written. Can include extensions like `.txt`, `.json`, etc.
            content (str): The content to be written into the file. Must be a string.

        Returns:
            bool: True if the file was successfully written.

        Logging:
            - Logs when the method begins and ends.
            - Logs confirmation of the file being written and its destination.
        """
        
        logging.info("Printing in write_file_to_gcs")
        
        """Writes (overwrites) a file in a GCS folder, supporting all file types."""
        blob = self.bucket.blob(f"{folder_name}/{file_name}")
        
        blob.upload_from_string(content)
        
        logging.info(f"File '{file_name}' written (overwritten) in '{folder_name}/'.")
        
        logging.info("End of write_file_to_gcs")
        
        return True
    
    def delete_file_from_gcs(self, folder_name, file_name):
        
        """
        Deletes a specified file from a Google Cloud Storage (GCS) folder.

        This method checks if the file exists in the given folder within the GCS bucket.
        If the file is found, it deletes the file (blob) from the bucket.
        If the file does not exist, it returns an appropriate error message.

        Args:
            folder_name (str): The folder (prefix path) in the GCS bucket where the file is located.
            file_name (str): The name of the file to delete.

        Returns:
            tuple:
                - bool: True if deletion was successful, False otherwise.
                - str: Message indicating success or error reason.

        Logging:
            - Logs when the method starts and ends.
            - Returns early if file does not exist, preventing unnecessary delete attempts.

        Notes:
            - Google Cloud Storage uses a flat namespace; folder structures are part of blob names.
            - The method relies on `self.file_exists()` to check for file existence before deletion.
        """
        
        logging.info("Printing in delete_file_from_gcs")
        
        """Deletes a file from a GCS folder."""
        if not self.file_exists(folder_name, file_name):
            return False, f"Error: File '{file_name}' does not exist in '{folder_name}/'."

        blob = self.bucket.blob(f"{folder_name}/{file_name}")
        blob.delete()
        
        logging.info("End of delete_file_from_gcs")
        return True, f"Success: File '{file_name}' deleted from '{folder_name}/'."

    def file_exists(self, folder_name, file_name):
        
        """
        Checks if a specific file exists within a folder in the Google Cloud Storage (GCS) bucket.

        Args:
            folder_name (str): The folder (prefix path) in the GCS bucket to check.
            file_name (str): The name of the file to check for existence.

        Returns:
            bool: True if the file exists in the specified folder, False otherwise.

        Notes:
            - Uses the `blob.exists()` method to check for the presence of the file.
            - GCS uses a flat namespace; folder structures are simulated via blob name prefixes.
        """
        
        blob = self.bucket.blob(f"{folder_name}/{file_name}")
        return blob.exists()

    def upload_file_to_gcs(self, folder_name, local_file_path, destination_file_name):
        
        """
        Uploads a local file to a specified folder within a Google Cloud Storage (GCS) bucket.

        Args:
            folder_name (str): The target folder (prefix) in the GCS bucket where the file will be uploaded.
            local_file_path (str): The full path of the local file to upload.
            destination_file_name (str): The filename to use for the uploaded file in GCS.

        Returns:
            tuple:
                - bool: True if upload succeeded, False if an error occurred.
                - str: The full GCS path of the uploaded file if successful,
                       otherwise an error message describing the failure.
        """
        try:
            blob = self.bucket.blob(f"{folder_name}/{destination_file_name}")
            blob.upload_from_filename(local_file_path)
            gcs_file_path = f"gs://{self.bucket.name}/{folder_name}/{destination_file_name}"
            
            
            logging.info(f"File '{local_file_path}' uploaded as '{destination_file_name}' in '{folder_name}/'.")
            return True, gcs_file_path
        except Exception as e:
            return False, f"Error uploading file: {str(e)}"
    
    def read_json_file_from_gcs(self, folder_name, file_name):
        
        """
        Reads a JSON file from a specified folder in Google Cloud Storage (GCS) and parses its content.

        Args:
            folder_name (str): The GCS folder (prefix) containing the JSON file.
            file_name (str): The name of the JSON file to read.

        Returns:
            tuple:
                - bool: True if the file was successfully read and parsed; False otherwise.
                - dict or str: Parsed JSON content as a Python dictionary if successful;
                               otherwise, an error message describing the failure.
        """
        success, file_data = self.read_file_from_gcs(folder_name, file_name)
        if not success:
            return False, f"Error: Unable to read file '{file_name}' in '{folder_name}/'."
        
        try:
            json_content = json.loads(file_data.get("text", ""))
            return True, json_content
        except json.JSONDecodeError:
            return False, f"Error: Invalid JSON format in '{file_name}'."
    
    def list_documents_in_gcs(self, folder_name):
        
        """
        Lists all document (file) names within a specified folder in a Google Cloud Storage (GCS) bucket.

        Args:
            folder_name (str): The folder (prefix) path inside the GCS bucket to list documents from.

        Returns:
            list of str: A list of document names (file names) found inside the specified folder.
                         The folder prefix is stripped from the returned names.

        Notes:
            - Assumes that `self.client` is a GCS client instance and `self.bucket_name` is the GCS bucket name.
            - Excludes the folder prefix itself from the returned list.
        """
        # Get the bucket
        bucket = self.bucket
        folder_prefix = f"{folder_name}/"
        
        # List the files (documents) in the folder (prefix) inside the bucket
        blobs = self.client.list_blobs(self.bucket_name, prefix=folder_prefix)

        # Store the document names by stripping the folder prefix
        documents = [blob.name[len(folder_prefix):] for blob in blobs if blob.name != folder_prefix]

        return documents
    
    def read_file_from_gcs(self, folder_name, file_name):
        
        """
        Reads a file from a specified folder in a Google Cloud Storage (GCS) bucket and
        attempts to determine whether the file is text or binary.

        Args:
            folder_name (str): The folder (prefix) path inside the GCS bucket.
            file_name (str): The name of the file to read within the folder.

        Returns:
            tuple:
                - bool: Success status (True if file exists and was read, False otherwise).
                - dict or str: 
                    If successful, returns a dictionary with keys:
                        - "text": Decoded UTF-8 string if the file is text, otherwise None.
                        - "bytes": Raw bytes content of the file.
                    If unsuccessful, returns an error message string.
        """
        
        if not self.file_exists(folder_name, file_name):
            return False, f"Error: File '{file_name}' not found in '{folder_name}/'."

        blob = self.bucket.blob(f"{folder_name}/{file_name}")
        file_content = blob.download_as_bytes()  # Always get bytes

        # Try decoding as UTF-8 (text files), else return raw bytes
        try:
            text_content = file_content.decode("utf-8")
            return True, {"text": text_content, "bytes": file_content}
        except UnicodeDecodeError:
            return True, {"text": None, "bytes": file_content}  # Binary file
        
    def read_txt_file_from_gcs(self, folder_name, file_name):
        
        """
        Reads a text file from a specified folder in a Google Cloud Storage (GCS) bucket.

        This method uses `read_file_from_gcs` to fetch the file content and verifies
        that the file is valid UTF-8 text. If the file is binary or cannot be decoded
        as text, an error message is returned.

        Args:
            folder_name (str): The folder (prefix) path inside the GCS bucket.
            file_name (str): The name of the text file to read.

        Returns:
            str: 
                - The UTF-8 decoded text content if the file is valid.
                - An error message string if the file cannot be read or is not valid text.
        """
        success, file_data = self.read_file_from_gcs(folder_name, file_name)
        if not success:
            return f"Error: Unable to read file '{file_name}' in '{folder_name}/'."
        
        if file_data.get("text") is None:
            
            logging.info(f"Error: '{file_name}' is not a valid text file.")
            return f"Error: '{file_name}' is not a valid text file."
        
        return file_data["text"]
    
    def extract_table_data_from_gcs(self, folder_name, file_name):
        
        """
        Extracts table data from an Excel (.xlsx) or CSV file stored in a Google Cloud Storage (GCS) folder.

        For Excel files:
            - Processes all visible sheets.
            - Skips sheets with insufficient rows or duplicate/blank column names.
            - Extracts the first 15 rows from each valid sheet and converts them to JSON lines format.
            - Aggregates JSON outputs from all sheets, prefixing each with the sheet name.

        For CSV files:
            - Reads the entire CSV file.
            - Skips the file if it contains duplicate or blank column names.
            - Extracts the first 15 rows and converts them to JSON lines format.

        Args:
            folder_name (str): The GCS folder (prefix) containing the file.
            file_name (str): The name of the Excel or CSV file.

        Returns:
            str: A string containing JSON lines representation of the extracted data 
                 prefixed with sheet names for Excel files, or an error message string 
                 if the file cannot be read or has invalid structure.
        """
        
        logging.info("Printing in extract_table_data_from_gcs")
        
        """Extracts table data from an Excel or CSV file stored in GCS."""
        success, file_data = self.read_file_from_gcs(folder_name, file_name)
        if not success:
            
            logging.info(f"Error: Unable to read file '{file_name}' in '{folder_name}/'.")
            return f"Error: Unable to read file '{file_name}' in '{folder_name}/'."

        file_bytes = file_data.get("bytes", None)
        if file_bytes is None:
            
            logging.info(f"Error: Could not retrieve file content.")
            return f"Error: Could not retrieve file content."

        # Read file as DataFrame
        file_ext = file_name.split(".")[-1].lower()
        data = ''

        if file_ext == "xlsx":
            workbook = load_workbook(filename=BytesIO(file_bytes))
            unhidden_sheets = [sheet.title for sheet in workbook.worksheets if sheet.sheet_state == 'visible']

            for sheet in unhidden_sheets:
                
                logging.info("Processing sheet: ", sheet)
                data += f"sheetname: {sheet}\n"
                active_sheet = workbook[sheet]
                sheet_data = [row for row in active_sheet.iter_rows(min_row=1, values_only=True)]

                if not sheet_data or len(sheet_data) < 2:
                    
                    logging.info(f"Skipping sheet '{sheet}' due to insufficient data.")
                    data += f"Skipping sheet '{sheet}' due to insufficient data.\n"
                    continue

                df = pd.DataFrame(data=sheet_data[1:], columns=sheet_data[0])

                # Check for duplicate or blank column names
                if df.columns.duplicated().any() or any(pd.isna(df.columns)) or any(df.columns == ''):
                    # print(f"Skipping sheet '{sheet}' due to duplicate or blank column names.")
                    logging.info(f"Skipping sheet '{sheet}' due to duplicate or blank column names.")
                    data += f"Skipping sheet '{sheet}' due to duplicate or blank column names.\n"
                    continue

                df = df.head(15)
                
                json_data = df.to_json(orient='records', lines=True)
                data += json_data + '\n'

        else:  # Handle CSV files
            data += "sheetname: sheet_1\n"
            df = pd.read_csv(BytesIO(file_bytes))

            # Check for duplicate or blank column names
            if df.columns.duplicated().any() or any(pd.isna(df.columns)) or any(df.columns == ''):
                
                logging.info(f"Skipping CSV file due to duplicate or blank column names.")
                return f"Skipping CSV file due to duplicate or blank column names."

            df = df.head(15)
            json_data = df.to_json(orient='records', lines=True)
            data += json_data + '\n'
        
        logging.info("End of extract_table_data_from_gcs")
        
        return data

    def read_docx_from_gcs(self, folder_name, file_name):
        
        logging.info("printing in read_docx_from_gcs")
        
        """
        Reads a DOCX file stored in a Google Cloud Storage (GCS) folder and extracts its full text content.

        Args:
            folder_name (str): The GCS folder (prefix) containing the DOCX file.
            file_name (str): The name of the DOCX file to read.

        Returns:
            str: The extracted text content of the DOCX file as a single string with paragraphs separated by newlines,
                 or an error message if the file cannot be read or its content is unavailable.
        """
        success, file_data = self.read_file_from_gcs(folder_name, file_name)
        if not success:
            
            logging.info(f"Error: Unable to read file '{file_name}' in '{folder_name}/'.")
            return f"Error: Unable to read file '{file_name}' in '{folder_name}/'."
        
        file_bytes = file_data.get("bytes", None)
        if file_bytes is None:
            
            logging.info(f"Error: Could not retrieve file content.")
            return f"Error: Could not retrieve file content."
        
        doc = Document(BytesIO(file_bytes))
        content = "\n".join([para.text for para in doc.paragraphs])
        
        logging.info("End of read_docx_from_gcs")
        return content
    
    def read_xml_from_gcs(self, folder_name, file_name):
        
        """
        Reads an XML file stored in a Google Cloud Storage (GCS) folder and returns its text content as a string.

        Args:
            folder_name (str): The GCS folder (prefix) containing the XML file.
            file_name (str): The name of the XML file to read.

        Returns:
            str: The raw text content of the XML file, or an error message if the file cannot be read.
        """
        success, file_data = self.read_file_from_gcs(folder_name, file_name)
        if not success:
            
            logging.info(f"Error: Unable to read file '{file_name}' in '{folder_name}/'.")
            return f"Error: Unable to read file '{file_name}' in '{folder_name}/'."
        
        xml_content=file_data['text']
        return xml_content
    import pandas as pd

    def read_pdf_from_gcs(self, folder_name, file_name, image=True):
        
        """
        Reads a PDF file from a Google Cloud Storage (GCS) folder and extracts its content.

        Args:
            folder_name (str): The GCS folder containing the PDF file.
            file_name (str): The name of the PDF file to read.
            image (bool, optional): If False, extracts text directly from the PDF pages.
                                    If True, converts pages to images and extracts text via OCR.
                                    Defaults to True.

        Returns:
            str or list:
                - If image=False: Returns the extracted text from all pages as a single string.
                - If image=True: Returns a string containing OCR-extracted text per page.
                - Returns an empty list on exception during image-based extraction.
                - Returns an error message string if the file cannot be read.
        """
        from PIL import Image
        success, file_data = self.read_file_from_gcs(folder_name, file_name)
        if not success:
            
            logging.info(f"Error: Unable to read file '{file_name}' in '{folder_name}/'.")
            return f"Error: Unable to read file '{file_name}' in '{folder_name}/'."
        
        # Download the PDF as bytes
        pdf_bytes= file_data['bytes']
        
        # Open the PDF from the byte data using PyMuPDF (fitz)
        pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")

        # Extract text from all pages
        if image==False:
            text = ""
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                text += page.get_text()
            return text  
        
        else:
            try:
                page_details={}
                for i, page in enumerate(pdf_document):
                    pix = page.get_pixmap(dpi=300)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    # Save image to bytes in PNG format
                    img_byte_arr = BytesIO()
                    img.save(img_byte_arr, format='PNG')
                    extracted_text=self.read_image_data(image_bytes=img_byte_arr.getvalue())
                    page_details[f"{i+1}"]=extracted_text
                    text='\n\n'.join(f"{k}: {v}" for k,v in page_details.items())
                
                return text
                
            except Exception as e:
                return []
    
    def read_image_from_gcs(self, folder_name, file_name):
        
        """
        Reads an image file (PNG, JPG, JPEG) from a Google Cloud Storage (GCS) folder and extracts its text content using OCR.

        Args:
            folder_name (str): The GCS folder containing the image file.
            file_name (str): The name of the image file to read.

        Returns:
            str: Extracted text content from the image, stripped of leading/trailing whitespace.
        """
        success, file_data = self.read_file_from_gcs(folder_name, file_name)
 
        # Download the Image as bytes
        image_bytes= file_data['bytes']
        extracted_text=self.read_image_data(image_bytes=image_bytes)
        return extracted_text.strip()  # Returns full structured ER details
    
    def encode_file_to_base64_from_gcs(self, folder_name, file_name):
        
        """
        Reads a file from a GCS folder and encodes its contents to a Base64 string.

        Args:
            folder_name (str): The GCS folder where the file is stored.
            file_name (str): The name of the file to encode.

        Returns:
            str: Base64-encoded string of the file contents if successful,
                 otherwise an error message string.
        """
        success, file_data = self.read_file_from_gcs(folder_name, file_name)
        if not success:
            
            logging.info(f"Error: Unable to read file '{file_name}' in '{folder_name}/'.")
            return f"Error: Unable to read file '{file_name}' in '{folder_name}/'."
        
        file_bytes = file_data.get("bytes", None)
        if file_bytes is None:
            
            logging.info(f"Error: Could not retrieve file content.")
            return f"Error: Could not retrieve file content."
        
        return base64.b64encode(file_bytes).decode("utf-8")
    
    def convert_xlsx_to_string_from_gcs(self, folder_name, file_name):
        
        """
        Reads an Excel (.xlsx) file from a GCS folder and converts its contents into a plain string.

        Args:
            folder_name (str): The GCS folder where the Excel file is stored.
            file_name (str): The name of the Excel file.

        Returns:
            tuple:
                - bool: True if reading and conversion succeed, False otherwise.
                - str: The string representation of the Excel sheet contents if successful,
                       otherwise an error message.
        """
        
        success, file_data = self.read_file_from_gcs(folder_name, file_name)
        if not success:
            return False, f"Error: Unable to read file '{file_name}' in '{folder_name}/'."
        
        file_bytes = file_data.get("bytes", None)
        if file_bytes is None:
            return False, f"Error: Could not retrieve file content."
        
        # Load workbook from bytes
        wb = openpyxl.load_workbook(BytesIO(file_bytes))
        sheet = wb.active
        
        # Convert sheet contents to string
        data_string = ""
        for row in sheet.iter_rows(values_only=True):
            row_string = "\t".join(str(cell) if cell is not None else "" for cell in row)
            data_string += row_string + "\n"
        
        return True, data_string

    def delete_file(self, file_path, stage):
        """
        Deletes a specified file from GCS if it exists.

        :param file_path: Path of the file in GCS.
        :param stage: Stage of the Explainable AI file.
        """
        file_path = file_path + f"/XAI/{stage}_ExplainableAI.txt"
        blob = self.bucket.blob(file_path)

        if blob.exists():
            blob.delete()
            print(f"Successfully deleted file: {file_path}")
            
            logging.info(f"Successfully deleted file: {file_path}")
            return True
        else:
            print(f"File {file_path} does not exist in bucket {self.bucket.name}.")
            
            logging.info(f"File {file_path} does not exist in bucket {self.bucket.name}.")
            return False

    def extract_excel_data_from_gcs(self,user_id,  file_name, sheet_name=None):
        
        folder_name=f'{user_id}/JIRA'
        
        """
        Extracts tabular data from an Excel (.xlsx) or CSV file stored in a GCS folder specific to a user.

        Args:
            user_id (str): The user identifier, used to build the GCS folder path (user_id/JIRA).
            file_name (str): The name of the file to read from GCS.
            sheet_name (str, optional): The specific sheet name to extract data from in an Excel file.
                                        Defaults to the first sheet if not provided.

        Returns:
            pd.DataFrame or str: A pandas DataFrame containing the extracted data if successful;
                                 otherwise, an error message string describing the failure.
        """
        success, file_data = self.read_file_from_gcs(folder_name, file_name)
        if not success:
            return f"Error: Unable to read file '{file_name}' in '{folder_name}/'."

        file_bytes = file_data.get("bytes", None)
        if file_bytes is None:
            return "Error: Could not retrieve file content."

        file_ext = file_name.split(".")[-1].lower()

        if file_ext == "xlsx":
            with BytesIO(file_bytes) as file_stream:
                workbook = load_workbook(filename=file_stream, data_only=True)

                if sheet_name is None:
                    sheet_name = workbook.sheetnames[0]  # Default to first sheet if not specified

                if sheet_name not in workbook.sheetnames:
                    return f"Error: Sheet '{sheet_name}' not found in the file."

                sheet = workbook[sheet_name]
                sheet_data = [row for row in sheet.iter_rows(min_row=1, values_only=True)]

                if not sheet_data or len(sheet_data) < 2:
                    return f"Error: Insufficient data in sheet '{sheet_name}'."

                df = pd.DataFrame(data=sheet_data[1:], columns=sheet_data[0])
                
        elif file_ext == "csv":
            df = pd.read_csv(BytesIO(file_bytes))

            if df.columns.duplicated().any() or any(pd.isna(df.columns)) or any(df.columns == ''):
                return "Error: Duplicate or blank column names in CSV file."

        else:
            return "Error: Unsupported file format. Only XLSX and CSV are supported."

        return df  # Returning the DataFrame directly
    
    def extract_text_from_image_with_gemini(self,image_data):
    
        print("printing in extract_text_from_image_with_gemini")
        
        """
        Extracts and logically structures text from an image stored in a GCS bucket using Gemini OCR.

        Args:
            gcs_image_path (str): GCS path to the image (e.g., 'gs://bucket-name/image.jpg').

        Returns:
            str: Extracted and formatted text.

        Raises:
            Exception: For Gemini model or GCS errors.
        """
        system_prompt = """
        You are an expert AI assistant specialized in high-accuracy Optical Character Recognition 
        (OCR) and intelligent document reconstruction. Your task is to extract all visible text and
          structural elements from a provided image, with a special focus on tables and layout 
          fidelity.
        Follow these instructions:
            1. Extract all text — including printed, handwritten, faint, rotated, or stylized text. 
            Do not skip any content.
            2. Reconstruct the content into a clear, logically organized, human-readable format. 
            This includes:
                a. Preserve and replicate tables:
                    - Detect all tables, including those with borders, no borders, or gridlines.
                    - Rebuild tables in markdown format, ensuring all rows, columns, headers, and 
                    cell contents are accurately represented.
                    - Handle merged or multiline cells with appropriate formatting 
                    (e.g., row spans, new lines, or indentations).
                b. Maintain document structure:
                    - Preserve headings, subheadings, paragraphs, bullet points, and numbering.
                    - Respect indentation, section breaks, and group related content based on 
                    layout cues.
                c. Correct OCR errors only when confidence is high. Avoid guessing unclear content.
                d. If visual clues suggest structure not labeled explicitly 
                (e.g., columns, headers, or sections), infer and apply structure.
            3. Output the final result as:
                - Clean, complete, and fully readable text.
                - All tables reconstructed accurately.
                - All layout elements preserved or clearly represented in textual format."""

        user_prompt = """Extract and format all text and structure from this image, including tables. " \
        "Return the result in human-readable form with tables in markdown."""
        try:
            gemini_image = Image.from_bytes(image_data)

            # Run Gemini OCR
            model = GenerativeModel(config.GEMINI_MODEL, system_instruction=system_prompt)
            generation_config = {
                "max_output_tokens": 8192,
                "temperature": 0.1,
                "top_p": 0.4,
            }
            
            
            response = model.generate_content(
                [gemini_image, user_prompt],
                generation_config=generation_config
            )
        
            
            return response.text.strip()

        except Exception as e:
            raise Exception(f"Error extracting text from image with Gemini: {type(e).__name__} - {e}")
    
    def read_file_gcs(self, folder_name, file_name):
        
        """
        Reads a file from Google Cloud Storage and returns its content as bytes.

        Args:
            folder_name (str): The GCS folder (prefix) where the file is located.
            file_name (str): The name of the file to read.

        Returns:
            bytes: The raw content of the file.
        """
        blob = self.bucket.blob(f"{folder_name}/{file_name}")
        # blob = self.bucket.blob(file_path)
        return blob.download_as_bytes()
    
    def pdf_to_images(self, folder_name, file_name, image_format="png"):
        
        """
        Convert each page of a PDF file stored in GCS to images and return the image bytes.

        Args:
            folder_name (str): The folder in GCS containing the PDF.
            file_name (str): The PDF file name.
            image_format (str): The desired image format (default: "png").

        Returns:
            tuple: (image_data, total_pages)
                - image_data (list of tuples): Each tuple contains (image_name, image_bytes).
                - total_pages (int): Total number of pages processed.
        """
        print("folder_name: ",folder_name)
        print("file_name: ",file_name)
        
        logging.info("folder_name: ",folder_name)
        logging.info("file_name: ",file_name)
        file_bytes = self.read_file_gcs(folder_name, file_name)
        print("\033[91mread_file_gcs in pdf_to_images\033[0m")
        
        logging.info("\033[91mread_file_gcs in pdf_to_images\033[0m")
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        image_data = []
        

        for i, page in enumerate(doc):
            total_page=i
            pix = page.get_pixmap(dpi=300)
            img_bytes = pix.tobytes(output=image_format)
            image_data.append((f"{file_name.rsplit('.', 1)[0]}_page_{i+1}.{image_format}", img_bytes))
            

        return image_data,total_page  # List of tuples (image_name, image_bytes)


    
    def extract_file_content(self,filepath: str) -> str:
        print("printing in extract_file_content")
        
        logging.info("printing in extract_file_content")
        """
        Extracts text content from a single supported file.
        
        Args:
            filepath (str): Path to the file.
        
        Returns:
            str: Extracted text content or error message.
        """
        folder_name = config.input_folder
        print("filename in extract_file_content: ",filepath)
        
        logging.info("filename in extract_file_content: ",filepath)
        file_extension = os.path.splitext(filepath)[1].lower()
        combined_text = ""
        total_page = ""
        try:
            if file_extension in [".doc", ".docx"]:
                print("\033[91mprocessing docx\033[0m")
                
                logging.info("\033[91mprocessing docx\033[0m")
                combined_text = self.read_docx_from_gcs(folder_name,filepath)
                return combined_text, total_page
            
            elif file_extension in ['.livedoc']:
                print("\033[91mprocessing doc and livedoc\033[0m")
                
                logging.info("\033[91mprocessing doc and livedoc\033[0m")
                combined_text = self.extract_text_from_image_with_gemini(filepath).decode(errors='ignore')
                return combined_text, total_page
            
            elif file_extension == '.pdf':
                print("\033[91mprocessing pdf\033[0m")
                
                logging.info("\033[91mprocessing pdf\033[0m")
                document_name = os.path.splitext(os.path.basename(filepath))[0]
                gcs_image_paths, total_page = self.pdf_to_images(folder_name, filepath)
                print("\033[91mprocessing pdf_to_images\033[0m")
                
                logging.info("\033[91mprocessing pdf_to_images\033[0m")
                extracted_sections = ''
                # print("gcs_image_paths.items(): ",gcs_image_paths[0])
                for image_name, image_bytes in gcs_image_paths:
                    print(f"Processing image: {image_name}")
                    
                    logging.info(f"Processing image: {image_name}")
                    try:
                        # image = Img.from_bytes(image_bytes)
                        print("\033[91mExtracting text from image\033[0m")
                        
                        logging.info("\033[91mExtracting text from image\033[0m")
                        text = self.extract_text_from_image_with_gemini(image_bytes)
                        extracted_sections += text.strip() + " "
                        time.sleep(1.5)
                    except Exception as e:
                        error_msg = f"Failed to process extraction of detail table for file '{image_name}': {e}"
                        print(error_msg)
                        return error_msg, total_page

                combined_text = f"Name of the document: {document_name}, and it contains information: {extracted_sections.strip()}"
                print("done for PDF")
                
                logging.info("done for PDF")
                return combined_text, total_page
            elif file_extension in [".csv", ".xlsx"]:
                combined_text = self.extract_table_data_from_gcs(folder_name, filepath)

            # Process XML
            elif file_extension in [".xml"]:
                combined_text = self.read_xml_from_gcs(folder_name, filepath) 

            elif file_extension in ['.ppt', '.pptx']:
                print("\033[91mprocessing ppt and pptx\033[0m")
                
                logging.info("\033[91mprocessing ppt and pptx\033[0m")
                prs = Presentation(filepath)
                combined_text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
                
                print("combined_text: ", combined_text)
                return combined_text, total_page

            elif file_extension in ['.jpeg', '.jpg', '.png']:
                print("\033[91mprocessing jpeg, hpg, png\033[0m")
                
                logging.info("\033[91mprocessing jpeg, hpg, png\033[0m")
                combined_text = self.read_image_from_gcs(folder_name,filepath)
                return combined_text, total_page

            elif file_extension == '.json':
                print("\033[91mprocessing json\033[0m")
                
                logging.info("\033[91mprocessing json\033[0m")
                with open(filepath, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                combined_text = json.dumps(data, indent=2)
                return combined_text, total_page
            
            elif file_extension == '.aspx':
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    soup = BeautifulSoup(f, 'html.parser')
                return soup.get_text()
            
            else:
                print("\033[91mprinting in else\033[0m")
                
                logging.info("\033[91mprinting in else\033[0m")
                return f"Unsupported file type: {file_extension}", total_page

        except Exception as e:
            return f"Error processing {filepath}: {type(e).__name__} - {e}", total_page


    def extract_all_from_folder(self, folder_path: str) -> pd.DataFrame:
        """
        Extracts text from all supported files in a GCS folder OR a single file.
        Also extracts and processes files inside ZIP archives.

        Args:
            folder_path (str): GCS path to a folder or single file (e.g., 'folder/subfolder/').

        Returns:
            pd.DataFrame: DataFrame with columns ['file', 'content', 'page_doc_type'].

        Raises:
            FileNotFoundError: If no files are found at the specified GCS path.
            Exception: If an unexpected error occurs during processing.
        """
        if not folder_path.endswith("/"):
            folder_path += "/"

        extracted_data = []

        try:
            blobs = list(self.bucket.list_blobs(prefix=folder_path))
            if not blobs:
                raise FileNotFoundError(f"No files found at GCS path: {folder_path}")

            print(f"Found {len(blobs)} blobs in path '{folder_path}'")
            
            logging.info(f"Found {len(blobs)} blobs in path '{folder_path}'")

            for blob in blobs:
                if blob.name.endswith('/'):
                    continue  # skip folders

                print(f"Processing blob: {blob.name}")
                
                logging.info(f"Processing blob: {blob.name}")
                ext = os.path.splitext(blob.name)[1].lower()
                rel_path = os.path.splitext(blob.name[len(folder_path):].lstrip('/'))[0]
                original_filename = os.path.basename(blob.name)
                base_file_name = os.path.splitext(original_filename)[0]

                try:
                    file_data = blob.download_as_bytes()

                    if ext == '.zip':
                        with tempfile.TemporaryDirectory() as temp_dir:
                            with zipfile.ZipFile(io.BytesIO(file_data)) as zip_ref:
                                zip_ref.extractall(temp_dir)
                                for zip_root, _, zip_files in os.walk(temp_dir):
                                    for zip_file in zip_files:
                                        zip_file_path = os.path.join(zip_root, zip_file)
                                        zip_rel_path = os.path.splitext(os.path.relpath(zip_file_path, temp_dir))[0]
                                        file_name = os.path.join(base_file_name, zip_rel_path)

                                        try:
                                            content, page_doc_type = self.extract_file_content(zip_file_path)
                                            extracted_data.append({
                                                "file": file_name,
                                                "content": content,
                                                "page_doc_type": page_doc_type
                                            })
                                        except Exception as e:
                                            print(f"Error extracting from ZIP file {file_name}: {e}")
                                            
                                            logging.info(f"Error extracting from ZIP file {file_name}: {e}")
                                            extracted_data.append({
                                                "file": file_name,
                                                "content": f"Error reading file: {type(e).__name__} - {e}",
                                                "page_doc_type": None
                                            })
                    else:
                        # Use original extension when creating temp file
                        suffix = os.path.splitext(blob.name)[1]
                        file_path = os.path.basename(blob.name)
                        base_file_name = os.path.splitext(os.path.basename(file_path))[0]
                        
                        print(f"(original: {blob.name})")
    
                        logging.info(f"(original: {blob.name})")

                        content, page_doc_type = self.extract_file_content(file_path)
                        extracted_data.append({
                            "file": base_file_name,
                            "content": content,
                            "page_doc_type": page_doc_type
                        })

                except Exception as e:
                    print(f"Error processing blob {blob.name}: {e}")
                    
                    logging.info(f"Error processing blob {blob.name}: {e}")
                    extracted_data.append({
                        "file": base_file_name,
                        "content": f"Error reading file: {type(e).__name__} - {e}",
                        "page_doc_type": None
                    })

            return pd.DataFrame(extracted_data)

        except Exception as e:
            raise Exception(f"Error during GCS path processing: {type(e).__name__} - {e}")